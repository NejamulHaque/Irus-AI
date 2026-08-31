import os
import io
import json
import csv
import re
import zipfile
import xml.etree.ElementTree as ET

import requests
import numpy as np
from werkzeug.utils import secure_filename
from PyPDF2 import PdfReader
from docx import Document as DocxDocument
from pptx import Presentation
from openpyxl import load_workbook

from app.models import db, Document, DocumentChunk

ALLOWED_EXTENSIONS = {
    # Documents & Presentations
    'pdf', 'docx', 'doc', 'pptx', 'ppt', 'rtf', 'odt', 'odp', 'epub',
    # Spreadsheets & Tabular Data
    'xlsx', 'xls', 'csv', 'tsv', 'ods',
    # Plain Text, Markdown & Documentation
    'txt', 'md', 'markdown', 'rst', 'log',
    # Structured Data & Configs
    'json', 'yaml', 'yml', 'xml', 'toml', 'ini', 'cfg', 'env',
    # Source Code & Scripts
    'py', 'js', 'jsx', 'ts', 'tsx', 'html', 'htm', 'css', 'scss',
    'c', 'cpp', 'h', 'hpp', 'java', 'go', 'rs', 'php', 'rb', 'sql', 'sh', 'bash', 'zsh'
}


def allowed_file(filename):
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS


# ---------- Specialized Format Extractors ----------

def _extract_pdf(stream):
    reader = PdfReader(stream)
    pages = []
    for i, page in enumerate(reader.pages, 1):
        txt = (page.extract_text() or "").strip()
        if txt:
            pages.append(f"[Page {i}]\n{txt}")
    return "\n\n".join(pages)


def _extract_docx(stream):
    try:
        doc = DocxDocument(stream)
        text_parts = []
        for para in doc.paragraphs:
            if para.text.strip():
                text_parts.append(para.text.strip())
        # Also extract all tables cleanly in markdown format
        for t_idx, table in enumerate(doc.tables, 1):
            rows = []
            for row in table.rows:
                row_cells = [c.text.strip().replace('\n', ' ') for c in row.cells]
                if any(row_cells):
                    rows.append(" | ".join(row_cells))
            if rows:
                text_parts.append(f"\n[Table {t_idx}]\n" + "\n".join(rows))
        return "\n\n".join(text_parts)
    except Exception:
        # Fallback for older .doc or XML-based documents
        stream.seek(0)
        return _extract_xml_from_zip(stream, 'word/document.xml')


def _extract_pptx(stream):
    try:
        prs = Presentation(stream)
        slides_text = []
        for i, slide in enumerate(prs.slides, 1):
            slide_content = []
            title = ""
            if slide.shapes.title and slide.shapes.title.text:
                title = f": {slide.shapes.title.text.strip()}"
            slide_content.append(f"[Slide {i}{title}]")

            for shape in slide.shapes:
                if shape.has_text_frame and shape != slide.shapes.title:
                    txt = shape.text.strip()
                    if txt:
                        slide_content.append(txt)
                elif shape.has_table:
                    t_rows = []
                    for row in shape.table.rows:
                        cells = [c.text.strip().replace('\n', ' ') for c in row.cells]
                        if any(cells):
                            t_rows.append(" | ".join(cells))
                    if t_rows:
                        slide_content.append("\n".join(t_rows))

            # Check for speaker notes
            try:
                if slide.has_notes_slide and slide.notes_slide.notes_text_frame.text:
                    notes = slide.notes_slide.notes_text_frame.text.strip()
                    if notes:
                        slide_content.append(f"[Speaker Notes]: {notes}")
            except Exception:
                pass

            slides_text.append("\n".join(slide_content))
        return "\n\n---\n\n".join(slides_text)
    except Exception:
        # Fallback XML parsing
        stream.seek(0)
        return _extract_xml_from_zip(stream, r'ppt/slides/slide\d+\.xml')


def _extract_xlsx(stream):
    wb = load_workbook(stream, data_only=True)
    sheets_text = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows_content = [f"[Sheet: {sheet_name}]"]
        row_count = 0
        for row in ws.iter_rows(values_only=True):
            if any(row):
                row_cells = [str(c).strip() if c is not None else "" for c in row]
                rows_content.append(" | ".join(row_cells))
                row_count += 1
                if row_count >= 1000:  # Prevent runaway huge sheets
                    rows_content.append(f"... [Truncated after {row_count} rows] ...")
                    break
        if len(rows_content) > 1:
            sheets_text.append("\n".join(rows_content))
    return "\n\n---\n\n".join(sheets_text)


def _extract_csv(stream, is_tsv=False):
    raw = stream.read()
    content = _decode_bytes(raw)
    reader = csv.reader(io.StringIO(content), delimiter='\t' if is_tsv else ',')
    rows = []
    for row in reader:
        if any(row):
            rows.append(" | ".join([c.strip() for c in row]))
    return "\n".join(rows)


def _extract_opendocument(stream):
    with zipfile.ZipFile(stream) as z:
        if 'content.xml' in z.namelist():
            tree = ET.fromstring(z.read('content.xml'))
            texts = [elem.text.strip() for elem in tree.iter() if elem.text and elem.text.strip()]
            return "\n\n".join(texts)
    return ""


def _extract_epub(stream):
    text_parts = []
    with zipfile.ZipFile(stream) as z:
        for name in z.namelist():
            if name.endswith(('.html', '.xhtml', '.htm')):
                raw = z.read(name)
                html = _decode_bytes(raw)
                # Strip HTML tags
                clean = re.sub(r'<style.*?</style>', '', html, flags=re.DOTALL)
                clean = re.sub(r'<script.*?</script>', '', clean, flags=re.DOTALL)
                clean = re.sub(r'<[^>]+>', ' ', clean)
                clean = re.sub(r'\s+', ' ', clean).strip()
                if clean:
                    text_parts.append(clean)
    return "\n\n".join(text_parts)


def _extract_rtf(stream):
    raw = stream.read()
    text = _decode_bytes(raw)
    text = re.sub(r'\\par[d]?', '\n', text)
    text = re.sub(r'\\tab', '\t', text)
    text = re.sub(r'\\[a-zA-Z0-9]+ ?', '', text)
    text = re.sub(r'[{}]', '', text)
    return "\n".join([line.strip() for line in text.splitlines() if line.strip()])


def _extract_html(stream):
    raw = stream.read()
    html = _decode_bytes(raw)
    clean = re.sub(r'<style.*?</style>', '', html, flags=re.DOTALL)
    clean = re.sub(r'<script.*?</script>', '', clean, flags=re.DOTALL)
    clean = re.sub(r'<br\s*/?>', '\n', clean)
    clean = re.sub(r'</p>', '\n\n', clean)
    clean = re.sub(r'</div>', '\n', clean)
    clean = re.sub(r'<[^>]+>', ' ', clean)
    return "\n".join([line.strip() for line in clean.splitlines() if line.strip()])


def _extract_structured_data(stream, ext):
    raw = stream.read()
    text = _decode_bytes(raw)
    if ext == 'json':
        try:
            parsed = json.loads(text)
            return f"[JSON Structure]\n" + json.dumps(parsed, indent=2)
        except Exception:
            return text
    return text


def _extract_plain_text(stream):
    raw = stream.read()
    return _decode_bytes(raw)


def _decode_bytes(b):
    for enc in ('utf-8', 'utf-16', 'latin-1', 'cp1252'):
        try:
            return b.decode(enc)
        except (UnicodeDecodeError, Exception):
            continue
    return b.decode('utf-8', errors='ignore')


def _extract_xml_from_zip(stream, pattern):
    try:
        with zipfile.ZipFile(stream) as z:
            texts = []
            for name in z.namelist():
                if re.search(pattern, name):
                    tree = ET.fromstring(z.read(name))
                    texts.extend([elem.text.strip() for elem in tree.iter() if elem.text and elem.text.strip()])
            return "\n\n".join(texts)
    except Exception:
        return ""


def extract_text(stream, filename):
    """Extract text from any supported document stream."""
    ext = filename.rsplit('.', 1)[1].lower() if '.' in filename else 'txt'
    try:
        if ext == 'pdf':
            return _extract_pdf(stream).strip()
        elif ext in ('docx', 'doc'):
            return _extract_docx(stream).strip()
        elif ext in ('pptx', 'ppt'):
            return _extract_pptx(stream).strip()
        elif ext in ('xlsx', 'xls'):
            return _extract_xlsx(stream).strip()
        elif ext in ('csv', 'tsv'):
            return _extract_csv(stream, is_tsv=(ext == 'tsv')).strip()
        elif ext in ('odt', 'ods', 'odp'):
            return _extract_opendocument(stream).strip()
        elif ext == 'epub':
            return _extract_epub(stream).strip()
        elif ext == 'rtf':
            return _extract_rtf(stream).strip()
        elif ext in ('html', 'htm'):
            return _extract_html(stream).strip()
        elif ext in ('json', 'yaml', 'yml', 'xml', 'toml', 'ini', 'cfg', 'env'):
            return _extract_structured_data(stream, ext).strip()
        else:
            return _extract_plain_text(stream).strip()
    except Exception as e:
        print(f"Error extracting text from {filename}: {e}")
        try:
            stream.seek(0)
            return _extract_plain_text(stream).strip()
        except Exception:
            return ""


def chunk_text(text, chunk_size=1000, overlap=200):
    chunks = []
    start = 0
    while start < len(text):
        chunks.append(text[start:start + chunk_size])
        start += (chunk_size - overlap)
    return [c for c in chunks if c.strip()]


# ---------- Semantic embeddings (free local Ollama) ----------

def _embedding_enabled():
    return os.getenv('EMBEDDING_PROVIDER', 'ollama').lower() != 'off'


def get_embeddings(texts):
    if not _embedding_enabled() or not texts:
        return None
    base = os.getenv('OLLAMA_BASE_URL', 'http://localhost:11434').rstrip('/')
    if base.endswith('/v1'):
        base = base[:-3]
    model = os.getenv('EMBEDDING_MODEL', 'nomic-embed-text')
    try:
        r = requests.post(f"{base}/api/embed", json={"model": model, "input": texts}, timeout=180)
        if r.ok:
            return r.json().get('embeddings')
        vecs = []
        for t in texts:
            r2 = requests.post(f"{base}/api/embeddings", json={"model": model, "prompt": t}, timeout=60)
            if not r2.ok:
                return None
            vecs.append(r2.json().get('embedding'))
        return vecs
    except Exception as e:
        print(f"--- DEBUG: embeddings unavailable ({e}). Using TF-IDF fallback. ---")
        return None


def cosine_similarity(a, b):
    a = np.asarray(a, dtype=np.float32)
    b = np.asarray(b, dtype=np.float32)
    denom = np.linalg.norm(a) * np.linalg.norm(b)
    return float(np.dot(a, b) / denom) if denom else 0.0


# ---------- Cloud-safe processing (NO disk writes) ----------

def process_and_save_document(file, user_id):
    if not allowed_file(file.filename):
        raise ValueError("File type not allowed.")

    filename = secure_filename(file.filename)
    stream = io.BytesIO(file.read())

    text = extract_text(stream, filename)
    if not text:
        raise ValueError("Could not extract text from file (it might be empty or a scanned image).")

    doc = Document(
        user_id=user_id,
        filename=filename,
        original_name=file.filename,
        file_path=''  # cloud-safe: knowledge lives in DB chunks, not on disk
    )
    db.session.add(doc)
    db.session.commit()

    chunks = chunk_text(text)
    vectors = get_embeddings(chunks)

    for i, chunk_content in enumerate(chunks):
        vec_json = None
        if vectors and i < len(vectors) and vectors[i]:
            vec_json = json.dumps(vectors[i])
        db.session.add(DocumentChunk(
            document_id=doc.id,
            content=chunk_content,
            chunk_index=i,
            embedding=vec_json
        ))

    db.session.commit()
    return doc


def get_relevant_context(query, document_id, top_k=3):
    chunks = DocumentChunk.query.filter_by(
        document_id=document_id
    ).order_by(DocumentChunk.chunk_index).all()

    if not chunks:
        return ""

    if len(chunks) <= top_k:
        return "\n\n---\n\n".join([c.content for c in chunks])

    vecs = [json.loads(c.embedding) if c.embedding else None for c in chunks]

    # Lazy re-index older documents
    if _embedding_enabled() and any(v is None for v in vecs):
        missing_idx = [i for i, v in enumerate(vecs) if v is None]
        new_vecs = get_embeddings([chunks[i].content for i in missing_idx])
        if new_vecs:
            for j, idx in enumerate(missing_idx):
                vecs[idx] = new_vecs[j]
                chunks[idx].embedding = json.dumps(new_vecs[j])
            try:
                db.session.commit()
            except Exception:
                db.session.rollback()

    # SEMANTIC SEARCH
    if all(v is not None for v in vecs):
        qvec = get_embeddings([query])
        if qvec and qvec[0]:
            scored = sorted(
                range(len(chunks)),
                key=lambda i: cosine_similarity(qvec[0], vecs[i]),
                reverse=True
            )
            top = sorted(scored[:top_k])
            return "\n\n---\n\n".join([chunks[i].content for i in top])

    # TF-IDF FALLBACK
    from sklearn.feature_extraction.text import TfidfVectorizer
    from sklearn.metrics.pairwise import cosine_similarity as tfidf_cosine

    chunk_texts = [c.content for c in chunks]
    vectorizer = TfidfVectorizer().fit(chunk_texts + [query])
    chunk_vectors = vectorizer.transform(chunk_texts)
    query_vector = vectorizer.transform([query])
    similarities = tfidf_cosine(query_vector, chunk_vectors).flatten()
    top_indices = similarities.argsort()[-top_k:][::-1]
    return "\n\n---\n\n".join([chunks[i].content for i in sorted(top_indices)])